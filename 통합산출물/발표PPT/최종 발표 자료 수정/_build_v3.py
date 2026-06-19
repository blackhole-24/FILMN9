# -*- coding: utf-8 -*-
"""FINSIGHT 발표 V3 — Appendix 디자인 톤으로 본편 재구성 + 밸류 복사 + Appendix/Q&A 첨부."""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ── 색 ──
NAVY=RGBColor(0x16,0x31,0x4F); GOLD=RGBColor(0xC7,0x92,0x2E); BLUE=RGBColor(0x2D,0x6C,0xB5)
INDIGO=RGBColor(0x4F,0x46,0xE5); GREEN=RGBColor(0x1E,0x8E,0x54); GRAY=RGBColor(0x64,0x72,0x7F)
INK=RGBColor(0x1F,0x2A,0x37); BAND=RGBColor(0xF7,0xFA,0xFD); WHITE=RGBColor(0xFF,0xFF,0xFF)
AMBER=RGBColor(0xFE,0xF3,0xC7); BROWN=RGBColor(0x92,0x40,0x0E); RED=RGBColor(0xB9,0x1C,0x1C)
LBLUE=RGBColor(0xEC,0xF1,0xF9); BORD=RGBColor(0xD9,0xE2,0xEF)
F="맑은 고딕"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BL=prs.slide_layouts[6]

def slide(navy_bg=False):
    s=prs.slides.add_slide(BL)
    if navy_bg:
        s.background.fill.solid(); s.background.fill.fore_color.rgb=NAVY
    return s
def rect(s,l,t,w,h,fill=None,rounded=False,line=None,lw=1.0):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                           Inches(l),Inches(t),Inches(w),Inches(h))
    if rounded:
        try: shp.adjustments[0]=0.08
        except: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    return shp
def txt(s,l,t,w,h,runs,size=14,color=INK,bold=False,italic=False,align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP,wrap=True,sp=None):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in (tf.margin_left,): pass
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs,str): runs=[(runs,{})]
    first=True
    for line in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align
        if sp is not None: p.space_after=Pt(sp)
        segs=line if isinstance(line,list) else [line]
        for seg in segs:
            t_,o=seg
            r=p.add_run(); r.text=t_
            r.font.size=Pt(o.get('size',size)); r.font.bold=o.get('bold',bold)
            r.font.italic=o.get('italic',italic); r.font.name=o.get('font',F)
            r.font.color.rgb=o.get('color',color)
    return tb
def frame(s,tag,title,sub=None,note=None,marker=None,title_size=34):
    rect(s,0.60,0.50,0.07,0.46,GOLD)
    txt(s,0.80,0.46,12.0,0.34,[[(tag,{})]],size=15,color=GOLD,bold=True)
    txt(s,0.60,0.84,12.4,0.92,[[(title,{})]],size=title_size,color=NAVY,bold=True)
    if sub: txt(s,0.62,1.74,12.2,0.5,[[(sub,{})]],size=16.5,color=GRAY)
    if note:
        rect(s,0.60,6.92,0.05,0.34,GOLD)
        txt(s,0.76,6.86,11.9,0.44,[[(note,{})]],size=12.5,color=GRAY,italic=True,anchor=MSO_ANCHOR.MIDDLE)
    if marker:
        txt(s,12.0,7.02,1.2,0.3,[[(marker,{})]],size=11,color=GRAY,align=PP_ALIGN.RIGHT)

def card(s,l,t,w,h,accent,head,headsize=14.5,headsub=None,body=None,fill=BAND,bar=True):
    rect(s,l,t,w,h,fill,rounded=True)
    if bar: rect(s,l,t,0.13,h,accent)
    yy=t+0.18
    txt(s,l+0.30,yy,w-0.5,0.4,[[(head,{})]],size=headsize,color=accent,bold=True)
    yy+=0.42
    if headsub:
        txt(s,l+0.30,yy,w-0.5,0.3,[[(headsub,{})]],size=10.5,color=GRAY); yy+=0.34
    if body:
        txt(s,l+0.30,yy,w-0.5,h-(yy-t)-0.1,[[ (seg if isinstance(seg,tuple) else (seg,{})) ] for seg in body],
            size=12,color=INK,sp=5)

# ════════ 1. 표지 ════════
s=slide(navy_bg=True)
txt(s,0.85,2.35,11,1.1,[[("FINSIGHT",{})]],size=56,color=WHITE,bold=True)
rect(s,0.90,3.55,3.2,0.06,GOLD)
txt(s,0.90,3.78,11,0.6,[[("AI 기반 기업 분석 자동화 서비스",{})]],size=23,color=RGBColor(0xF3,0xC9,0x7A),bold=True)
txt(s,0.92,4.62,11,0.5,[[("공시(DART)·시세(KRX) 기반 · 어려운 정보를 쉽게 · 모든 수치에 출처",{})]],size=15.5,color=RGBColor(0xC2,0xCC,0xDC))
txt(s,0.92,6.55,11,0.4,[[("KPMG AI Lab   ·   FILMN9 Inc.   ·   2026.06.20",{})]],size=13,color=RGBColor(0x8B,0x97,0xAC))

# ════════ 2. 목차 ════════
s=slide(); frame(s,"CONTENTS","발표 목차")
items=[("01","기획 배경 · 페르소나","문제와 해결방안 — 왜·누구를 위해"),
       ("02","개발 과정 & 시스템 설계","WBS · 아키텍처 · 파이프라인 · ERD · 기술 스택"),
       ("03","프로젝트 팀 역할 분담","PM · 기업분석 · 밸류에이션 · AI 챗봇"),
       ("04","비즈니스 모델","개발 예산 · 수익 구조(Free/Pro) · BEP"),
       ("05","밸류에이션 설명","멀티플 · 피어 · DCF · WACC"),
       ("06","AWS 배포 버전 시연","실서비스 라이브 접속"),
       ("07","Appendix","기술 상세 부록"),
       ("08","Q & A","질의응답")]
for i,(n,tt,sb) in enumerate(items):
    col=i//4; row=i%4
    x=0.7+col*6.25; y=2.45+row*1.12
    rect(s,x,y,5.95,0.92,BAND,rounded=True); rect(s,x,y,0.10,0.92,GOLD if col else NAVY)
    txt(s,x+0.28,y+0.20,1.0,0.5,[[(n,{})]],size=23,color=(GOLD if col else NAVY),bold=True)
    txt(s,x+1.35,y+0.17,4.4,0.4,[[(tt,{})]],size=15.5,color=NAVY,bold=True)
    txt(s,x+1.35,y+0.55,4.5,0.3,[[(sb,{})]],size=10.5,color=GRAY)

# ════════ 3. 기획 배경 ════════
s=slide(); frame(s,"01 · 기획 배경","왜 만들었나 — 정보의 비대칭을 줄인다",
    "전문가의 정보를, 모두에게. 개인 투자자의 정보 격차를 좁힙니다.",
    "정보 비대칭은 줄이고, 신뢰는 출처로. 거짓 없는 팩트만 — NO-MOCK 원칙.","")
card(s,0.6,2.55,6.0,3.9,RED,"현실 (AS-IS)",15,"개인 투자자가 마주하는 벽",
     ["정확한 기업 정보에 접근·해석이 어렵다",
      "커뮤니티·유튜브 정보는 출처가 불명확하다",
      "재무·공시는 용어가 어렵고 어디서 볼지 모른다",
      "지금 주가가 비싼지 싼지 판단이 안 된다"])
card(s,6.85,2.55,6.0,3.9,GREEN,"FINSIGHT (TO-BE)",15,"우리가 푸는 방식",
     ["공시(DART) 기반 정확한 정보를 쉬운 말로 제공",
      "모든 수치에 출처 — 누구나 즉시 팩트체크",
      "한 회사를 5분 안에 파악하도록 구조화",
      "적정가·해석까지, 스스로 판단하도록 지원"])

# ════════ 4. 페르소나 ════════
s=slide(); frame(s,"01 · 페르소나","타깃 사용자 — 직장인 이준호 (33)",
    "주식을 시작했지만, 무엇을 어떻게 봐야 할지 막막한 개인 투자자.",
    "목표 — 이준호가 한 종목을 5분 안에 '뭐 하는 회사·돈은 버는지·비싼지·믿을지' 판단하게 한다.","")
card(s,0.6,2.55,4.3,3.9,BLUE,"페르소나",15,None,
     ["• 33세 · 직장인 · 투자 경력 1년",
      "• 월 50만 원 적립식 투자",
      "• 점심시간·퇴근 후 짧게 확인",
      "• 손실은 두렵지만 공부할 시간은 부족",
      "• '근거 있는 정보'를 원함"])
card(s,5.15,2.55,7.7,3.9,RED,"문제점 (AS-IS)",15,"이준호가 겪는 불편",
     ["종목을 추천받아도 왜 좋은지 스스로 판단하지 못한다",
      "재무제표·공시는 어렵고, 어디서 봐야 할지 모른다",
      "정보의 출처가 불명확해 신뢰하기 어렵다",
      "지금 주가가 적정한지(비싼지 싼지) 감이 오지 않는다"])

# ════════ 5. WBS ════════
s=slide(); frame(s,"02 · 개발 과정","WBS — 킥오프에서 최종 발표까지",
    "5개 스프린트 · 멘토링 4회 반영 · 스프린트 4부터 파트별 역할 분담.",
    "S4부터 기업분석 파트(재무·브리핑·관리자)와 밸류에이션 파트(엔진·4-Way·챗봇)로 분담.","02")
rows=[("킥오프 (4/22~24)","프로젝트 헌장 · 페르소나 · WBS 초안","PM · 공통"),
      ("S1 (4/27~28)","기획서 · 요구사항정의서 · 협업환경 구축","PM · 공통"),
      ("S2 (4/29~5/9)","개발환경 · 데이터 수집 · 밸류 엔진(WACC·DCF·멀티플)","밸류에이션 파트"),
      ("S3 (5/12~22)","기업분석 · AI 히스토리 브리핑 · Next.js UI · 1차 배포","기업분석 / 밸류에이션"),
      ("S4 (5/26~6/5)","재무 정확성(3개년) · 브리핑 전수 · 4-Way 밸류 · 관리자","기업분석 / 밸류에이션"),
      ("S5 (6/8~19)","DB 클라우드 이관 · 계열사 전종목 · 메인 리디자인 · 챗봇 GPU · AWS+HTTPS · QA","기업분석 / 밸류에이션"),
      ("최종 발표 (6/20)","최종 발표 · 라이브 시연","전원")]
cw=[3.0,7.0,2.7]; x0=0.6; y0=2.50; rh=0.47
# header
hx=x0
for j,htxt in enumerate(["단계 (기간)","핵심 작업","담당 파트"]):
    rect(s,hx,y0,cw[j],rh,NAVY); txt(s,hx+0.14,y0,cw[j]-0.2,rh,[[(htxt,{})]],size=12.5,color=WHITE,bold=True,anchor=MSO_ANCHOR.MIDDLE)
    hx+=cw[j]
for i,row in enumerate(rows):
    y=y0+rh*(i+1); hx=x0; fill=WHITE if i%2==0 else LBLUE
    for j,cell in enumerate(row):
        rect(s,hx,y,cw[j],rh,fill,line=BORD,lw=0.5)
        col=NAVY if j==0 else (BLUE if (j==2 and "/" in cell) else INK)
        bold=(j==0) or (j==2)
        txt(s,hx+0.14,y,cw[j]-0.2,rh,[[(cell,{'size':11.5})]],size=11.5,color=col,bold=bold,anchor=MSO_ANCHOR.MIDDLE)
        hx+=cw[j]

# ════════ 7. 데이터 파이프라인 ════════
s=slide(); frame(s,"02 · 데이터 파이프라인","데이터 파이프라인 — 수집에서 화면까지",
    "공식 출처 → 수집·정제 → 저장 → 가공 → 화면. 모든 단계에서 출처를 보존합니다.",
    "주가는 매 거래일 16시 자동 갱신 · 재무·공시는 분기 갱신 · 모든 수치는 출처 추적 가능.","02")
stages=[("①  수집","DART · KRX · ECOS\nKOFIA · 네이버",BLUE),
        ("②  정제·표준화","XBRL 파싱 · 단위 정규화\n결측 검증",BLUE),
        ("③  저장","RDS · MongoDB\nChromaDB · S3",INDIGO),
        ("④  가공","밸류 엔진 · RAG\nSankey · 검증",INDIGO),
        ("⑤  제공","FastAPI →\nNext.js 화면",GREEN)]
sw=2.32; gap=0.18; x=0.6; y=3.0; h=2.4
for i,(h1,b,acc) in enumerate(stages):
    rect(s,x,y,sw,h,BAND,rounded=True); rect(s,x,y,sw,0.13,acc)
    txt(s,x+0.2,y+0.30,sw-0.35,0.45,[[(h1,{})]],size=16,color=acc,bold=True)
    txt(s,x+0.2,y+0.95,sw-0.35,1.3,[[(ln,{'size':12})] for ln in b.split("\n")],size=12,color=INK,sp=3)
    if i<4: txt(s,x+sw-0.02,y+0.85,0.34,0.5,[[("▶",{})]],size=15,color=GOLD,align=PP_ALIGN.CENTER)
    x+=sw+gap

# ════════ 8. ERD ════════
s=slide(); frame(s,"02 · 데이터 구조","ERD",
    "데이터 성격에 맞춰 4개 데이터베이스에 분산 저장합니다 (폴리글랏).",
    "관계형은 RDS · 문서형은 Mongo · 벡터는 Chroma · 파일은 S3 — 목적에 맞는 저장소 선택.","02")
dbs=[("RDS PostgreSQL","관계형 · 16 테이블",NAVY,
      ["company_info · financials","financial_detail · ohlcv","shareholders · executives","disclosures · valuation_summary"]),
     ("MongoDB Atlas","문서형",GREEN,
      ["histories","(AI 히스토리 브리핑)","회사별 연혁·사업요약","price_factors 등"]),
     ("ChromaDB","벡터 (임베딩)",INDIGO,
      ["사업보고서 RAG 청크","BGE-M3 1024차원","챗봇·브리핑 검색","리랭커 연계"]),
     ("AWS S3","파일 객체",BLUE,
      ["Sankey 손익흐름도","밸류 결과 JSON","계열사 시각화 이미지","정적 산출물"])]
cw2=3.02; x=0.6; y=2.7; h=3.5
for i,(nm,sub,acc,tbls) in enumerate(dbs):
    rect(s,x,y,cw2,h,BAND,rounded=True); rect(s,x,y,cw2,0.13,acc)
    txt(s,x+0.22,y+0.28,cw2-0.4,0.4,[[(nm,{})]],size=14,color=acc,bold=True)
    txt(s,x+0.22,y+0.70,cw2-0.4,0.3,[[(sub,{})]],size=10.5,color=GRAY)
    txt(s,x+0.22,y+1.15,cw2-0.4,h-1.3,[[(t,{'size':11.5})] for t in tbls],size=11.5,color=INK,sp=6)
    x+=cw2+0.13

# ════════ 9. 기술 스택 ════════
s=slide(); frame(s,"02 · 기술 스택","기술 스택",
    "프론트엔드 · 백엔드 · 데이터 · AI · 인프라 전 계층을 직접 구축했습니다.",
    "AI는 '찾기'만 맡고 판단은 정해진 규칙이 — 환각을 막는 구조로 설계.","02")
stack=[("프론트엔드",BLUE,"Next.js 16 · React 19 · TypeScript · Tailwind CSS"),
       ("백엔드",NAVY,"FastAPI · Python · Uvicorn · systemd"),
       ("데이터베이스",INDIGO,"PostgreSQL(RDS) · MongoDB Atlas · ChromaDB · SQLite"),
       ("AI · RAG",BROWN,"OpenAI gpt-5-mini · BGE-M3 임베딩 · bge-reranker"),
       ("밸류 엔진",GREEN,"DCF · EPV · 멀티플(4종) · 피어 자동선정"),
       ("인프라",BLUE,"AWS EC2·RDS·S3 · nginx · Let's Encrypt(HTTPS)")]
cw3=6.05; rh3=1.18; x0=0.6; y0=2.6
for i,(nm,acc,b) in enumerate(stack):
    col=i%2; rowi=i//2
    x=x0+col*6.25; y=y0+rowi*1.28
    rect(s,x,y,cw3,rh3,BAND,rounded=True); rect(s,x,y,0.13,rh3,acc)
    txt(s,x+0.30,y+0.18,cw3-0.5,0.4,[[(nm,{})]],size=14.5,color=acc,bold=True)
    txt(s,x+0.30,y+0.62,cw3-0.5,0.45,[[(b,{})]],size=12.5,color=INK)

# ════════ 10. 팀 구성 ════════
s=slide(); frame(s,"03 · 팀 구성","팀 구성 · 역할 분담",
    "파트별 책임 영역으로 나눠 개발하고, 통합·배포·QA는 공통으로 진행했습니다.",
    "각자 맡은 파트를 책임 개발 → 통합·배포·검증은 공통 영역에서 함께.","03")
team=[("이형주","PM · 공통 + 기업분석",NAVY,
       ["프로젝트 헌장 · 기획 · 요구사항","WBS · QA · 발표 총괄","재무 · 주가차트 · BS/IS · Sankey","전자공시 · 메인 · 통합검색","관리자 · AWS 배포 · 전수검수"]),
      ("고휘주","기업분석 (콘텐츠)",GREEN,
       ["AI 히스토리 브리핑 (RAG)","최신 뉴스","계열사 시각화","4-Phase 신뢰도 검증","MongoDB Atlas 적재"]),
      ("유태상","밸류에이션 + AI 챗봇",INDIGO,
       ["밸류 엔진 (DCF·멀티플·WACC)","밸류에이션 탭 전체","관리자 (밸류 운영·테스트)","AI 챗봇 RAG","피어 자동선정"])]
cw4=4.06; x=0.6; y=2.6; h=3.8
for nm,role,acc,duties in team:
    rect(s,x,y,cw4,h,BAND,rounded=True); rect(s,x,y,cw4,0.13,acc)
    txt(s,x+0.28,y+0.28,cw4-0.5,0.5,[[(nm,{})]],size=20,color=acc,bold=True)
    txt(s,x+0.28,y+0.82,cw4-0.5,0.35,[[(role,{})]],size=12.5,color=GRAY,bold=True)
    txt(s,x+0.28,y+1.30,cw4-0.5,h-1.5,[[("· "+d,{'size':12})] for d in duties],size=12,color=INK,sp=7)
    x+=cw4+0.19

# ════════ 11. 비즈니스 모델 ════════
s=slide(); frame(s,"04 · 비즈니스 모델","비즈니스 모델 — 개발 예산 · 수익 구조",
    "실제 결제 내역 기준. 단가 등 추정치는 '(가정)'으로 명시합니다 (NO-MOCK).",
    "AI 운영비를 로컬 GPU로 최소화 + AWS 크레딧 활용 → 실 현금지출 최소화.","04")
card(s,0.6,2.55,6.0,3.9,NAVY,"개발 예산 (원가)",15,"인건비 + 도구·인프라",
     ["인건비  ·  맨먼스 약 6 MM (단가 가정)",
      "Claude Code Max  ·  ₩187,000 / 월",
      "ChatGPT Plus  ·  약 ₩27,000 / 월",
      "OpenAI API(gpt-5-mini)  ·  사용량 비례",
      "AWS  ·  프리·크레딧 커버 (실부담 ₩0)",
      "도메인  ·  연 약 ₩15,000 (가정)"])
card(s,6.85,2.55,6.0,3.9,GREEN,"수익 모델 · 손익분기",15,"구독 기반",
     ["Free  ·  기본 분석 · 검색 · 히스토리 브리핑",
      "Pro(구독)  ·  심층 밸류 · 알림 · API (가격 가정)",
      "추가  ·  B2B 데이터 · 리포트 (확장)",
      "",
      "BEP  ·  Pro 구독자 N명 시 손익분기 (가정)",
      "고정비 작음 → 낮은 손익분기점"])

# ════════ 14. AWS 배포 시연 (divider) ════════
# (밸류 12·13, 아키텍처 6, 마무리 15, Appendix, QnA는 복사 단계에서 순서대로 삽입)
prs.save("_v3_built.pptx")
print("본편(복사전) 저장 완료:", len(prs.slides.__iter__.__self__._sldIdLst), "장")
