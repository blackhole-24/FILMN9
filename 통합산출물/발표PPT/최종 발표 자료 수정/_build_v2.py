# -*- coding: utf-8 -*-
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

BASE = "FINSIGHT_발표_최종_v1.pptx"
QNA  = "../FINSIGHT_QnA.pptx"
OUT  = "FINSIGHT_발표_최종_v2.pptx"

prs = Presentation(BASE)
qna = Presentation(QNA)

# ── cross-file 슬라이드 복사 (이미지 포함) ──
def blank_layout(p):
    for L in p.slide_layouts:
        if 'blank' in (L.name or '').lower(): return L
    return p.slide_layouts[-1]
def copy_slide(dst, src):
    ns = dst.slides.add_slide(blank_layout(dst))
    for sp in list(ns.shapes): sp._element.getparent().remove(sp._element)
    for sh in src.shapes: ns.shapes._spTree.append(copy.deepcopy(sh._element))
    rid = {}
    for r_id, rel in src.part.rels.items():
        if "image" in rel.reltype: rid[r_id] = ns.part.relate_to(rel.target_part, rel.reltype)
    if rid:
        for el in ns.shapes._spTree.iter():
            for a in (qn('r:embed'), qn('r:link')):
                v = el.get(a)
                if v in rid: el.set(a, rid[v])
    return ns

# QnA 5장 추가(append)
for s in qna.slides: copy_slide(prs, s)

sl = list(prs.slides)  # 0-based

# ── 헬퍼 ──
def set_text(shape, txt):
    p = shape.text_frame.paragraphs[0]
    if p.runs:
        p.runs[0].text = txt
        for r in p.runs[1:]: r._r.getparent().remove(r._r)
    else:
        p.add_run().text = txt
def by_name(slide, name):
    for sh in slide.shapes:
        if sh.name == name: return sh
    return None
def replace_all(slide, a, b):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if a in r.text: r.text = r.text.replace(a, b)

# ── 슬라이드 2 (목차) : 02~08 글자 교체 + 09·10칸 삭제 ──
s2 = sl[1]
NEW = {
 'Text 13':"개발 과정 & 시스템 설계",   'Text 14':"WBS·아키텍처·파이프라인·ERD·기술 스택",
 'Text 18':"프로젝트 팀 역할분담",       'Text 19':"PM · 기업분석 · 밸류에이션 · AI 챗봇",
 'Text 23':"비즈니스 모델",             'Text 24':"개발 예산 · 수익 구조 (Free/Pro · BEP)",
 'Text 28':"밸류에이션 설명",           'Text 29':"멀티플 · 피어 · DCF · WACC",
 'Text 33':"AWS 배포 버전 시연",        'Text 34':"기업분석·밸류·챗봇·관리자 (실접속)",
 'Text 38':"Appendix",                 'Text 39':"기술 상세 부록 — 공식·RAG·검증·인프라",
 'Text 43':"Q & A",                    'Text 44':"질의응답",
}
for nm, tx in NEW.items():
    sh = by_name(s2, nm)
    if sh: set_text(sh, tx)
DEL = {'Shape 45','Shape 46','Text 47','Text 48','Text 49','Shape 50','Shape 51','Text 52','Text 53','Text 54'}
for sh in list(s2.shapes):
    if sh.name in DEL: sh._element.getparent().remove(sh._element)

# ── 슬라이드 19·20 : 기업개요 → 기업분석 ──
replace_all(sl[18], "기업개요", "기업분석")
replace_all(sl[19], "기업개요", "기업분석")

# ── 슬라이드 23 (마무리) : 면책 문구 추가 (Appendix 전 마지막) ──
s23 = sl[22]
DISC = ("※ 본 자료와 서비스의 모든 분석·적정주가는 공개데이터(DART·KRX·한국은행·KOFIA 등)에 근거한 "
        "투자 참고 정보이며, 특정 종목의 매수·매도를 권유하지 않습니다. 투자 판단과 책임은 이용자 본인에게 있습니다.")
tb = s23.shapes.add_textbox(Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.5))
tb.name = "Disclaimer"
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run(); r.text = DISC
f = r.font; f.size = Pt(8); f.italic = True; f.name = "맑은 고딕"; f.color.rgb = RGBColor(0x6B,0x70,0x88)

# ── 슬라이드 재배열 (0-based) : 7·8 삭제, 목차순 + QnA ──
# v1: 표지0 목차1 | 기획3~6=2,3,4,5 | WBS9=8 아키13~16=12,13,14,15 | 팀10=9 | 비용11,12=10,11
#     | 밸류17,18=16,17 | 라이브22=21 관리자19,20,21=18,19,20 | 마무리23=22 | Appendix24~38=23~37 | QnA=38~42
order = [0,1, 2,3,4,5, 8,12,13,14,15, 9, 10,11, 16,17, 21,18,19,20, 22, *range(23,38), 38,39,40,41,42]
lst = prs.slides._sldIdLst
ids = list(lst)
for el in ids: lst.remove(el)
for i in order: lst.append(ids[i])

prs.save(OUT)
print("저장:", OUT, "· 최종", len(list(prs.slides)), "장 (목표 41)")
