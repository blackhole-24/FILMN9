# -*- coding: utf-8 -*-
import copy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

NAVY=RGBColor(0x0E,0x27,0x42); GOLD=RGBColor(0xB5,0x87,0x2B); WHITE=RGBColor(0xFF,0xFF,0xFF)
GRAYL=RGBColor(0xAF,0xB9,0xCC); CREAM=RGBColor(0xE6,0xC5,0x8A); F="맑은 고딕"
APX=r"C:\Users\Admin\Downloads\FINSIGHT_Appendix.pptx"; V2="FINSIGHT_발표_최종_v2.pptx"; QNA="../FINSIGHT_QnA.pptx"

prs=Presentation("_v4_built.pptx"); apx=Presentation(APX); v2=Presentation(V2); qna=Presentation(QNA)
def blank(p):
    for L in p.slide_layouts:
        if 'blank' in (L.name or '').lower(): return L
    return p.slide_layouts[-1]
def copy_slide(dst,src):
    ns=dst.slides.add_slide(blank(dst))
    for sp in list(ns.shapes): sp._element.getparent().remove(sp._element)
    cs=src._element.find(qn('p:cSld')); bg=cs.find(qn('p:bg')) if cs is not None else None
    if bg is not None: ns._element.find(qn('p:cSld')).insert(0,copy.deepcopy(bg))
    for sh in src.shapes: ns.shapes._spTree.append(copy.deepcopy(sh._element))
    rid={}
    for r_id,rel in src.part.rels.items():
        if "image" in rel.reltype: rid[r_id]=ns.part.relate_to(rel.target_part,rel.reltype)
    if rid:
        for el in ns.shapes._spTree.iter():
            for a in (qn('r:embed'),qn('r:link')):
                if el.get(a) in rid: el.set(a,rid[el.get(a)])
    return ns
def replace_exact(slide,old,new):
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip()==old:
            p=sh.text_frame.paragraphs[0]
            if p.runs:
                p.runs[0].text=new
                for r in p.runs[1:]: r._r.getparent().remove(r._r)

# 10: AWS 디바이더
d=prs.slides.add_slide(blank(prs)); d.background.fill.solid(); d.background.fill.fore_color.rgb=NAVY
tb=d.shapes.add_textbox(Inches(0),Inches(2.9),Inches(13.33),Inches(1.4)); p=tb.text_frame.paragraphs[0]
p.alignment=PP_ALIGN.CENTER; r=p.add_run(); r.text="AWS 배포 버전 시연"; r.font.size=Pt(48); r.font.bold=True; r.font.name=F; r.font.color.rgb=WHITE
ln=d.shapes.add_shape(MSO_SHAPE.RECTANGLE,Inches(5.67),Inches(4.45),Inches(2.0),Inches(0.06)); ln.fill.solid(); ln.fill.fore_color.rgb=GOLD; ln.line.fill.background(); ln.shadow.inherit=False
tb2=d.shapes.add_textbox(Inches(0),Inches(4.7),Inches(13.33),Inches(0.5)); p2=tb2.text_frame.paragraphs[0]
p2.alignment=PP_ALIGN.CENTER; r2=p2.add_run(); r2.text="실서비스 라이브 접속  ·  https://43.203.94.124.nip.io"; r2.font.size=Pt(16); r2.font.name=F; r2.font.color.rgb=CREAM

# 11·12: 밸류 (V2 #15,16)
copy_slide(prs,v2.slides[14]); copy_slide(prs,v2.slides[15])
# 13: 마무리 (Appendix #19 + 면책)
fin=copy_slide(prs,apx.slides[18]); replace_exact(fin,"부록 S","")
DISC=("※ 본 자료와 서비스의 모든 분석·적정주가는 공개데이터(DART·KRX·한국은행·KOFIA 등)에 근거한 투자 참고 정보이며, 특정 종목의 매수·매도를 권유하지 않습니다. 투자 판단과 책임은 이용자 본인에게 있습니다.")
tb=fin.shapes.add_textbox(Inches(0.7),Inches(6.95),Inches(12.0),Inches(0.45)); rr=tb.text_frame.paragraphs[0].add_run()
rr.text=DISC; rr.font.size=Pt(8); rr.font.italic=True; rr.font.name=F; rr.font.color.rgb=GRAYL
# 14~31: Appendix #1~18 / 32~36: QnA #1~5
for i in range(0,18): copy_slide(prs,apx.slides[i])
for i in range(0,5): copy_slide(prs,qna.slides[i])

order=[0,1,2,3,4,5,6,7,8,9, 11,12, 10, 13] + list(range(14,32)) + list(range(32,37))
lst=prs.slides._sldIdLst; ids=list(lst)
for el in ids: lst.remove(el)
for i in order: lst.append(ids[i])
prs.save("FINSIGHT_발표_최종_v4.pptx")
print("V4 저장:",len(list(prs.slides)),"장 (목표 37)")
