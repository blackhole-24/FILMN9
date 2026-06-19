# -*- coding: utf-8 -*-
"""FINSIGHT V4 본편 — 더 무거운 비즈니스 톤 + 목차세로 + 기획·페르소나 통합 + 다이어그램 이미지 + 비즈 1p."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ── 더 무거운 비즈니스 팔레트 (딥네이비·쿨톤) ──
NAVY=RGBColor(0x0E,0x27,0x42); NAVY2=RGBColor(0x1B,0x35,0x55); GOLD=RGBColor(0xB5,0x87,0x2B)
BLUE=RGBColor(0x2A,0x5E,0x9E); INDIGO=RGBColor(0x3B,0x37,0xB0); GREEN=RGBColor(0x1A,0x6E,0x47)
GRAY=RGBColor(0x5A,0x65,0x73); INK=RGBColor(0x1A,0x24,0x30); BAND=RGBColor(0xED,0xF1,0xF7)
WHITE=RGBColor(0xFF,0xFF,0xFF); AMBER=RGBColor(0xF6,0xE7,0xC3); BROWN=RGBColor(0x7A,0x52,0x12)
RED=RGBColor(0xA8,0x28,0x28); LBLUE=RGBColor(0xE6,0xEC,0xF4); BORD=RGBColor(0xCF,0xD8,0xE6)
F="맑은 고딕"
DIAG=os.path.join(os.path.dirname(__file__),"_diag")

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BL=prs.slide_layouts[6]
def slide(navy=False):
    s=prs.slides.add_slide(BL)
    if navy: s.background.fill.solid(); s.background.fill.fore_color.rgb=NAVY
    return s
def rect(s,l,t,w,h,fill=None,rounded=False,line=None,lw=1.0):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    if rounded:
        try: shp.adjustments[0]=0.06
        except: pass
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False; return shp
def txt(s,l,t,w,h,lines,size=14,color=INK,bold=False,italic=False,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,sp=None,wrap=True):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    if isinstance(lines,str): lines=[[(lines,{})]]
    first=True
    for ln in lines:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align
        if sp is not None: p.space_after=Pt(sp)
        for seg in (ln if isinstance(ln,list) else [ln]):
            t_,o=seg; r=p.add_run(); r.text=t_
            r.font.size=Pt(o.get('size',size)); r.font.bold=o.get('bold',bold); r.font.italic=o.get('italic',italic)
            r.font.name=F; r.font.color.rgb=o.get('color',color)
    return tb
def frame(s,tag,title,sub=None,note=None,marker=None,tsize=33):
    rect(s,0.60,0.50,0.07,0.46,GOLD)
    txt(s,0.80,0.46,12,0.34,[[(tag,{})]],15,GOLD,True)
    txt(s,0.60,0.84,12.4,0.9,[[(title,{})]],tsize,NAVY,True)
    if sub: txt(s,0.62,1.72,12.2,0.5,[[(sub,{})]],16,GRAY)
    if note:
        rect(s,0.60,6.92,0.05,0.34,GOLD)
        txt(s,0.76,6.86,11.9,0.44,[[(note,{})]],12.5,GRAY,italic=True,anchor=MSO_ANCHOR.MIDDLE)
    if marker: txt(s,12.0,7.02,1.2,0.3,[[(marker,{})]],11,GRAY,align=PP_ALIGN.RIGHT)
def img(s,path,top,height):
    from PIL import Image
    iw,ih=Image.open(path).size; ratio=iw/ih; w=height*ratio
    x=(13.333-w)/2
    s.shapes.add_picture(path,Inches(x),Inches(top),height=Inches(height))

# ════ 1. 표지 ════
s=slide(navy=True)
txt(s,0.85,2.35,11,1.1,[[("FINSIGHT",{})]],56,WHITE,True)
rect(s,0.90,3.55,3.2,0.06,GOLD)
txt(s,0.90,3.78,11,0.6,[[("AI 기반 기업 분석 · 가치 평가 정보 지원 서비스",{})]],22,RGBColor(0xE6,0xC5,0x8A),True)
txt(s,0.92,4.62,11,0.5,[[("공시(DART)·시세(KRX) 기반 · 어려운 정보를 쉽게 · 모든 수치에 출처",{})]],15.5,RGBColor(0xB9,0xC4,0xD6))
txt(s,0.92,6.55,11,0.4,[[("KPMG AI Lab   ·   FILMN9 Inc.   ·   2026.06.20",{})]],13,RGBColor(0x82,0x8F,0xA6))

# ════ 2. 목차 (좌측 큰 타이틀 / 우측 세로 1~8) ════
s=slide()
rect(s,0.60,0.85,0.07,0.46,GOLD)
txt(s,0.80,0.82,5,0.34,[[("CONTENTS",{})]],15,GOLD,True)
txt(s,0.60,1.25,5.6,1.0,[[("발표 목차",{})]],40,NAVY,True)
txt(s,0.64,2.6,5.4,2.5,[[("FINSIGHT가 무엇을, 왜,",{'size':17,'color':GRAY})],
                         [("어떻게 만들었는지 —",{'size':17,'color':GRAY})],
                         [("여덟 단계로 설명드립니다.",{'size':17,'color':GRAY})]],17,GRAY,sp=4)
toc=["기획 배경 · 페르소나","개발 과정 & 시스템 설계","프로젝트 팀 역할 분담","비즈니스 모델",
     "밸류에이션 설명","AWS 배포 버전 시연","Appendix","Q & A"]
x=6.7; y=0.95; rh=0.74; gap=0.05
for i,t_ in enumerate(toc):
    yy=y+i*(rh+gap)
    rect(s,x,yy,6.0,rh,BAND,rounded=True); rect(s,x,yy,0.09,rh,NAVY if i<6 else GOLD)
    txt(s,x+0.28,yy,0.85,rh,[[("%02d"%(i+1),{})]],19,(NAVY if i<6 else GOLD),True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,x+1.15,yy,4.7,rh,[[(t_,{})]],15.5,NAVY,True,anchor=MSO_ANCHOR.MIDDLE)

# ════ 3. 기획 배경 + 페르소나 (통합) ════
s=slide()
rect(s,0.60,0.50,0.07,0.46,GOLD); txt(s,0.80,0.46,12,0.34,[[("01 · 기획 배경 · 문제",{})]],15,GOLD,True)
txt(s,0.60,0.84,12,0.8,[[("기획 배경",{})]],33,NAVY,True)
txt(s,0.62,1.66,12.4,0.95,
    [[("방대하고 다양한 투자 정보의 바다에서, 거짓 없는 신뢰 정보만을 스스로 학습하고 옳고 그름을 판단해 위험을 헷지합니다.",{'size':15.5,'color':INK})],
     [("정보의 비대칭을 조금이나마 줄여 주는 — AI 기반 기업 분석 및 가치 평가 정보 지원 서비스입니다.",{'size':15.5,'color':INK})]],sp=3)
# 페르소나 띠
rect(s,0.60,2.78,12.13,0.62,NAVY,rounded=True)
txt(s,0.85,2.78,1.0,0.62,[[("👤",{'size':18,'color':WHITE})]],18,WHITE,anchor=MSO_ANCHOR.MIDDLE)
txt(s,1.45,2.78,11.0,0.62,[[("페르소나 · 일반 투자자  ",{'size':14.5,'color':WHITE,'bold':True}),
    ("— 주식은 시작했지만, 무엇을 어떻게 봐야 할지 막막한 개인 투자자",{'size':13.5,'color':RGBColor(0xC9,0xD3,0xE2)})]],anchor=MSO_ANCHOR.MIDDLE)
# 문제점 4박스
probs=[("“정보가 흩어져 있다”","공시는 DART, 시세는 증권앱,\n뉴스는 포털 — 매번 창을 옮긴다",RED),
       ("“용어가 너무 어렵다”","PER·DCF·WACC… 비전공자에겐\n외국어. 결국 ‘느낌’으로 산다",GOLD),
       ("“출처를 못 믿겠다”","유튜브·단톡방 ‘카더라’.\n근거를 확인할 길이 없다",BLUE),
       ("“비싼지 싼지 모르겠다”","적정가를 가늠할 기준이 없어\n사도 되는지 판단이 안 선다",INDIGO)]
bw=2.95; x=0.6; y=3.62; h=2.7
for tt,bd,acc in probs:
    rect(s,x,y,bw,h,BAND,rounded=True); rect(s,x,y,bw,0.12,acc)
    txt(s,x+0.24,y+0.30,bw-0.45,0.8,[[(tt,{})]],15,acc,True)
    txt(s,x+0.24,y+1.15,bw-0.45,h-1.3,[[(ln,{'size':11.5})] for ln in bd.split("\n")],11.5,INK,sp=3)
    x+=bw+0.13
rect(s,0.60,6.92,0.05,0.34,GOLD)
txt(s,0.76,6.86,11.9,0.44,[[("결국 문제는 ‘정보 비대칭’ — 이 문제들을 FINSIGHT가 어떤 정보로 푸는지 지금부터 화면으로 보여드립니다.",{})]],12.5,GRAY,italic=True,anchor=MSO_ANCHOR.MIDDLE)

# ════ 4. WBS (S2 기업분석 추가·담당 통일) ════
s=slide(); frame(s,"02 · 개발 과정","WBS — 킥오프에서 최종 발표까지",
    "5개 스프린트 · 멘토링 4회 반영 · 스프린트 2부터 파트별 역할 분담.",
    "기업분석 파트(재무·브리핑·관리자)와 밸류에이션 파트(엔진·4-Way·챗봇)가 각 스프린트를 병행 진행.","02")
rows=[("킥오프 (4/22~24)","프로젝트 헌장 · 페르소나 · WBS 초안","PM · 공통"),
      ("S1 (4/27~28)","기획서 · 요구사항정의서 · 협업환경 구축","PM · 공통"),
      ("S2 (4/29~5/9)","개발환경 구축 · 기초데이터 수집·정제(DART·KRX) · 밸류 엔진(WACC·DCF·멀티플)","기업분석 / 밸류에이션"),
      ("S3 (5/12~22)","기업분석 화면 · AI 히스토리 브리핑 · Next.js UI · 1차 배포","기업분석 / 밸류에이션"),
      ("S4 (5/26~6/5)","재무 정확성(3개년) · 브리핑 전수 · 4-Way 밸류 · 관리자","기업분석 / 밸류에이션"),
      ("S5 (6/8~19)","DB 클라우드 이관 · 계열사 전종목 · 메인 리디자인 · 챗봇 GPU · AWS+HTTPS · QA","기업분석 / 밸류에이션"),
      ("최종 발표 (6/20)","최종 발표 · 라이브 시연","전원")]
cw=[2.9,7.1,2.7]; x0=0.6; y0=2.50; rh=0.47
hx=x0
for j,h_ in enumerate(["단계 (기간)","핵심 작업","담당 파트"]):
    rect(s,hx,y0,cw[j],rh,NAVY); txt(s,hx+0.14,y0,cw[j]-0.2,rh,[[(h_,{})]],12.5,WHITE,True,anchor=MSO_ANCHOR.MIDDLE); hx+=cw[j]
for i,row in enumerate(rows):
    y=y0+rh*(i+1); hx=x0; fill=WHITE if i%2==0 else LBLUE
    for j,cell in enumerate(row):
        rect(s,hx,y,cw[j],rh,fill,line=BORD,lw=0.5)
        col=NAVY if j==0 else (BLUE if (j==2 and "/" in cell) else INK); bold=(j==0)or(j==2)
        txt(s,hx+0.14,y,cw[j]-0.2,rh,[[(cell,{'size':11})]],11,col,bold,anchor=MSO_ANCHOR.MIDDLE); hx+=cw[j]

# ════ 5·6·7. 아키텍처 / 파이프라인 / ERD (이미지) ════
for tag,title,fn,note,mk in [
    ("02 · 시스템 아키텍처","시스템 아키텍처","s2_0.png","AWS 클라우드 위에서 웹·API(24/7)와 GPU 챗봇(온디맨드)을 분리 운영 — 4개 저장소로 폴리글랏 구성.","02"),
    ("02 · 데이터 파이프라인","데이터 파이프라인","s3_0.png","수집 → 정규화·생성 → 검증(NO-MOCK) → 적재(4저장소) → 서빙. 배치·실시간 이원화.","02"),
    ("02 · 데이터 구조","ERD","s4_0.png","관계형 핵심 16테이블(운영 13)은 RDS, 문서·벡터·파일은 Mongo·Chroma·S3로 분리.","02")]:
    s=slide(); rect(s,0.60,0.50,0.07,0.46,GOLD); txt(s,0.80,0.46,12,0.34,[[(tag,{})]],15,GOLD,True)
    txt(s,0.60,0.84,12,0.7,[[(title,{})]],28,NAVY,True)
    img(s,os.path.join(DIAG,fn),1.78,4.95)
    rect(s,0.60,6.92,0.05,0.34,GOLD); txt(s,0.76,6.86,11.9,0.44,[[(note,{})]],12,GRAY,italic=True,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,12.0,7.02,1.2,0.3,[[(mk,{})]],11,GRAY,align=PP_ALIGN.RIGHT)

# ════ 8. 기술 스택 ════
s=slide(); frame(s,"02 · 기술 스택","기술 스택",
    "프론트엔드 · 백엔드 · 데이터 · AI · 인프라 전 계층을 직접 구축했습니다.",
    "AI는 '찾기'만 맡고 판단은 정해진 규칙이 — 환각을 막는 구조로 설계.","02")
stack=[("프론트엔드",BLUE,"Next.js 16 · React 19 · TypeScript · Tailwind CSS"),
       ("백엔드",NAVY,"FastAPI · Python · Uvicorn · systemd"),
       ("데이터베이스",INDIGO,"PostgreSQL(RDS) · MongoDB Atlas · ChromaDB · SQLite"),
       ("AI · RAG",BROWN,"OpenAI gpt-5-mini · BGE-M3 임베딩 · bge-reranker · 온톨로지"),
       ("밸류 엔진",GREEN,"DCF · EPV · 멀티플(4종) · 피어 자동선정"),
       ("인프라",NAVY2,"AWS EC2·RDS·S3 · nginx · Let's Encrypt(HTTPS)")]
cw3=6.05; rh3=1.18; x0=0.6; y0=2.6
for i,(nm,acc,b) in enumerate(stack):
    col=i%2; rowi=i//2; x=x0+col*6.25; y=y0+rowi*1.28
    rect(s,x,y,cw3,rh3,BAND,rounded=True); rect(s,x,y,0.13,rh3,acc)
    txt(s,x+0.30,y+0.18,cw3-0.5,0.4,[[(nm,{})]],14.5,acc,True)
    txt(s,x+0.30,y+0.62,cw3-0.5,0.45,[[(b,{})]],12.5,INK)

# ════ 9. 팀 구성 ════
s=slide(); frame(s,"03 · 팀 구성","팀 구성 · 역할 분담",
    "파트별 책임 영역으로 나눠 개발하고, 통합·배포·QA는 공통으로 진행했습니다.",
    "각자 맡은 파트를 책임 개발 → 통합·배포·검증은 공통 영역에서 함께.","03")
team=[("이형주","PM · 공통 + 기업분석",NAVY,
       ["프로젝트 헌장 · 기획 · 요구사항","WBS · QA · 발표 총괄","재무 · 주가차트 · BS/IS · Sankey","전자공시 · 메인 · 통합검색","관리자 · AWS 배포 · 전수검수"]),
      ("고휘주","기업분석 (콘텐츠)",GREEN,
       ["AI 히스토리 브리핑 (RAG)","최신 뉴스","계열사 시각화","4-Phase 신뢰도 검증","MongoDB Atlas 적재"]),
      ("유태상","밸류에이션 + AI 챗봇",INDIGO,
       ["밸류 엔진 (DCF·멀티플·WACC)","밸류에이션 탭 전체","관리자 (밸류 운영·테스트)","AI 챗봇 RAG","피어 자동선정"])]
cw4=4.06; x=0.6; y=2.6; h=3.8
for nm,role,acc,duties in team:
    rect(s,x,y,cw4,h,BAND,rounded=True); rect(s,x,y,cw4,0.13,acc)
    txt(s,x+0.28,y+0.28,cw4-0.5,0.5,[[(nm,{})]],20,acc,True)
    txt(s,x+0.28,y+0.82,cw4-0.5,0.35,[[(role,{})]],12.5,GRAY,True)
    txt(s,x+0.28,y+1.30,cw4-0.5,h-1.5,[[("· "+d,{'size':12})] for d in duties],12,INK,sp=7)
    x+=cw4+0.19

# ════ 10. 비즈니스 모델 (1p · 총액·BEP·유료기능) ════
s=slide(); frame(s,"04 · 비즈니스 모델","비즈니스 모델 — 예산 · 수익 · 손익분기",
    "실제 결제 기준. 단가·구독료 등 추정치는 '(가정)'으로 명시합니다 (NO-MOCK).",
    "AI 운영비를 로컬 GPU로 최소화 + AWS 크레딧 활용 → 실 현금지출 최소화.","04")
# 좌: 개발 예산 + 총액
card=rect(s,0.6,2.5,6.0,4.0,BAND,rounded=True); rect(s,0.6,2.5,0.13,4.0,NAVY)
txt(s,0.9,2.68,5.6,0.4,[[("개발 예산 (원가)",{})]],15,NAVY,True)
txt(s,0.9,3.12,5.6,2.5,
    [[("· 인건비  맨먼스 약 6 MM (단가 가정)",{'size':12})],
     [("· Claude Code Max  ₩187,000 / 월",{'size':12})],
     [("· ChatGPT Plus  약 ₩27,000 / 월",{'size':12})],
     [("· OpenAI API(gpt-5-mini)  사용량 비례",{'size':12})],
     [("· AWS  프리·크레딧 커버 (실부담 ₩0)",{'size':12})],
     [("· 도메인  연 약 ₩15,000 (가정)",{'size':12})]],sp=7)
rect(s,0.9,5.55,5.4,0.78,NAVY,rounded=True)
txt(s,1.15,5.55,5.0,0.78,[[("총 개발원가  ",{'size':13,'color':WHITE}),("약 2,450만원",{'size':18,'color':RGBColor(0xF1,0xCB,0x7A),'bold':True}),(" (인건비 단가 가정)",{'size':11,'color':RGBColor(0xC9,0xD3,0xE2)})]],anchor=MSO_ANCHOR.MIDDLE)
# 우: 수익 + BEP
rect(s,6.85,2.5,6.0,4.0,BAND,rounded=True); rect(s,6.85,2.5,0.13,4.0,GREEN)
txt(s,7.15,2.68,5.6,0.4,[[("수익 모델 · 손익분기",{})]],15,GREEN,True)
txt(s,7.15,3.12,5.6,2.5,
    [[("Free  ",{'size':12,'bold':True}),("기본 분석·검색·히스토리 브리핑",{'size':12})],
     [("Pro 유료화  ",{'size':12,'bold':True,'color':GREEN}),("심층 밸류(DCF 시나리오)·적정가 알림·",{'size':12})],
     [("        관심종목 리포트·AI 챗봇 무제한·API",{'size':12})],
     [("· 가격(가정)  Pro 월 9,900원",{'size':12})],
     [("· 운영비(가정)  서버·API 약 월 20만원",{'size':12})]],sp=6)
rect(s,7.15,5.45,5.4,0.92,GREEN,rounded=True)
txt(s,7.4,5.45,5.0,0.92,
    [[("손익분기  ",{'size':12,'color':WHITE}),("약 21명 구독",{'size':15,'color':RGBColor(0xCF,0xF0,0xDC),'bold':True})],
     [("구독 300명 시  개발원가 ",{'size':12,'color':WHITE}),("약 9개월 회수",{'size':15,'color':RGBColor(0xCF,0xF0,0xDC),'bold':True}),(" (가정)",{'size':10,'color':RGBColor(0xCF,0xF0,0xDC)})]],sp=2,anchor=MSO_ANCHOR.MIDDLE)

prs.save("_v4_built.pptx")
print("V4 본편 저장:", len(list(prs.slides)), "장 (10 기대)")
