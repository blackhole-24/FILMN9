# -*- coding: utf-8 -*-
"""창B PM 슬라이드 생성 — WBS 표 / 시스템아키텍처 / 데이터파이프라인 / ERD / 기술스택.
다이어그램은 통합산출물/diagrams/png/*.png (SVG→PNG 변환본) 삽입.
실행: python 통합산출물/발표PPT/_gen_pm_slides.py → FINSIGHT_PM슬라이드_창B.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

ROOT = Path(r"C:\Users\Admin\FILMN9\통합산출물")
PNG = ROOT / "diagrams" / "png"
OUT = ROOT / "발표PPT" / "FINSIGHT_PM슬라이드_창B.pptx"

NAVY  = RGBColor(0x1F, 0x3A, 0x5F)
BLUE  = RGBColor(0x2F, 0x55, 0x97)
INK   = RGBColor(0x0F, 0x17, 0x2A)
GRAY  = RGBColor(0x64, 0x74, 0x8B)
LINE  = RGBColor(0xD9, 0xDF, 0xE7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DONE  = RGBColor(0x37, 0x56, 0x23); DONE_BG = RGBColor(0xE2, 0xEF, 0xDA)
PROG  = RGBColor(0x80, 0x60, 0x00); PROG_BG = RGBColor(0xFF, 0xF2, 0xCC)
PLAN  = RGBColor(0x80, 0x80, 0x80); PLAN_BG = RGBColor(0xF2, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_title(slide, kicker, title):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.30), Inches(12.2), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = kicker
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = BLUE
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = title
    r2.font.size = Pt(28); r2.font.bold = True; r2.font.color.rgb = NAVY
    return tb


def add_image_fit(slide, png, top=Inches(1.45), maxw=Inches(12.3), maxh=Inches(5.7)):
    from PIL import Image
    w, h = Image.open(png).size
    ar = w / h
    tw, th = maxw, Emu(int(maxw / ar))
    if th > maxh:
        th, tw = maxh, Emu(int(maxh * ar))
    left = Emu(int((SW - tw) / 2))
    slide.shapes.add_picture(str(png), left, top, width=tw, height=th)


# ───────── Slide 1 : WBS 요약 표 ─────────
s = prs.slides.add_slide(BLANK)
add_title(s, "개발 여정 · WBS", "킥오프(4/22) → 최종 발표(6/20) · 5 스프린트")
rows = [
    ("단계 (기간)", "핵심 Task", "담당 파트", "상태"),
    ("킥오프 (4/22~4/24)", "프로젝트 헌장 · 페르소나 · WBS 초안", "PM·공통", "완료"),
    ("스프린트 1 (4/27~4/28)", "기획서 · 요구사항정의서 · 협업환경 구축", "PM·공통", "완료"),
    ("스프린트 2 (4/29~5/9)", "개발환경 · 데이터 수집 · 밸류엔진(WACC·DCF·멀티플)", "밸류에이션", "완료"),
    ("스프린트 3 (5/12~5/22)", "기업개요 · AI 히스토리 브리핑 · Next.js UI · 1차 배포", "기업개요·AI챗봇·밸류", "완료"),
    ("스프린트 4 (5/26~6/5)", "재무 정확성(3개년) · 브리핑 전수화 · 4-Way 밸류 · 관리자", "전 파트", "완료"),
    ("스프린트 5 (6/8~6/19)", "DB 클라우드 이관 · 계열사 전종목 · 메인 리디자인 · 챗봇 GPU · AWS+HTTPS · QA", "전 파트", "완료/진행"),
    ("최종 발표 (6/20)", "발표 PPT · 스크립트 · 리허설 · 라이브 시연", "전 파트", "예정"),
]
tbl_w = Inches(12.3); tbl_h = Inches(4.9)
gf = s.shapes.add_table(len(rows), 4, Inches(0.5), Inches(1.55), tbl_w, tbl_h)
tb = gf.table
tb.columns[0].width = Inches(2.5); tb.columns[1].width = Inches(6.6)
tb.columns[2].width = Inches(2.0); tb.columns[3].width = Inches(1.2)
for ci in range(4):
    tb.cell(0, ci).fill.solid(); tb.cell(0, ci).fill.fore_color.rgb = BLUE
for ri, row in enumerate(rows):
    for ci, val in enumerate(row):
        cell = tb.cell(ri, ci); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.1); cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
        para = cell.text_frame.paragraphs[0]; para.alignment = PP_ALIGN.LEFT
        run = para.add_run(); run.text = val
        run.font.size = Pt(11.5 if ri else 12)
        if ri == 0:
            run.font.bold = True; run.font.color.rgb = WHITE
        else:
            run.font.color.rgb = INK
            if ci == 0: run.font.bold = True; run.font.color.rgb = BLUE
            if ci == 3:
                para.alignment = PP_ALIGN.CENTER; run.font.bold = True
                st = val
                if st == "완료": cell.fill.solid(); cell.fill.fore_color.rgb = DONE_BG; run.font.color.rgb = DONE
                elif st == "예정": cell.fill.solid(); cell.fill.fore_color.rgb = PLAN_BG; run.font.color.rgb = PLAN
                else: cell.fill.solid(); cell.fill.fore_color.rgb = PROG_BG; run.font.color.rgb = PROG
            else:
                if ri % 2 == 0: cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF7,0xF9,0xFC)
                else: cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
for ri in range(len(rows)):
    tb.rows[ri].height = Inches(0.42 if ri else 0.4)
foot = s.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.5))
fp = foot.text_frame.paragraphs[0]; fr = fp.add_run()
fr.text = "‘리스크 큰 것부터 증명하고, 정확히 만들고, 세상에 올린다’ — 5개 스프린트 · 멘토링 4회 피드백 반영 · 담당: PM·공통 / 기업개요 / 밸류에이션 / AI 챗봇"
fr.font.size = Pt(11); fr.font.color.rgb = GRAY


# ───────── Slide 2~4 : 다이어그램 이미지 ─────────
diagrams = [
    ("시스템 아키텍처", "AWS 위에서 실제로 돌아가는 서비스", "system_architecture_0619.png"),
    ("데이터 파이프라인", "수집 → 정규화·생성 → 검증(NO-MOCK) → 적재 → 서빙", "data_pipeline.png"),
    ("데이터 구조 · ERD", "폴리글랏 4-DB · 종목코드로 묶인 16테이블", "erd.png"),
]
for kicker, title, fn in diagrams:
    s = prs.slides.add_slide(BLANK)
    add_title(s, kicker, title)
    p = PNG / fn
    if p.exists():
        add_image_fit(s, p)
    else:
        tb = s.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        tb.text_frame.paragraphs[0].add_run().text = f"[이미지 누락: {fn}]"


# ───────── Slide 5 : 기술 스택 ─────────
s = prs.slides.add_slide(BLANK)
add_title(s, "기술 스택", "무엇으로 만들었나 — 계층별 선택")
stack = [
    ("프론트엔드", ["Next.js 16 · React 19 (SSR·standalone)", "TypeScript · Tailwind CSS", "recharts 차트·시각화"]),
    ("백엔드 · API", ["FastAPI (Python) · REST :8090 비동기", "psycopg3 · pymongo (DB 드라이버)", "systemd · nginx (기동·프록시)"]),
    ("데이터 · 저장", ["PostgreSQL (RDS) · 16테이블", "MongoDB · ChromaDB (문서·벡터)", "AWS S3 (파일 서빙)"]),
    ("AI · RAG", ["BGE-M3 1024d · 리랭커 (로컬 GPU)", "OpenAI gpt-5-mini (답변·진단)", "자체 RAG 파이프라인 (LangChain 미사용)"]),
    ("인프라 · 운영", ["AWS EC2(GPU) · RDS · S3", "HTTPS · Let's Encrypt", ".env 런타임 로드 (시크릿 0 커밋)"]),
    ("개발 · 협업", ["Claude Code · Design · v0.dev", "GitHub (feature + PR)", "Notion · Google Drive"]),
]
cols, rows_n = 3, 2
gx, gy = Inches(0.5), Inches(1.65)
cw, ch = Inches(4.05), Inches(2.45)
gapx, gapy = Inches(0.18), Inches(0.22)
for i, (head, items) in enumerate(stack):
    r, c = divmod(i, cols)
    x = Emu(int(gx) + i % cols * (int(cw) + int(gapx)))
    y = Emu(int(gy) + r * (int(ch) + int(gapy)))
    box = s.shapes.add_shape(1, x, y, cw, ch)  # rectangle
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xF7, 0xF9, 0xFC)
    box.line.color.rgb = LINE; box.line.width = Pt(1)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.14)
    hp = tf.paragraphs[0]; hr = hp.add_run(); hr.text = head
    hr.font.size = Pt(15); hr.font.bold = True; hr.font.color.rgb = NAVY
    for it in items:
        bp = tf.add_paragraph(); br = bp.add_run(); br.text = "· " + it
        br.font.size = Pt(11.5); br.font.color.rgb = INK; bp.space_before = Pt(4)
foot = s.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.5))
fr = foot.text_frame.paragraphs[0].add_run()
fr.text = "AI는 자체 파이프라인으로 제어 — LangChain 미사용·경량화 / 모든 시크릿은 .env 런타임 로드(하드코딩·커밋 0)"
fr.font.size = Pt(11); fr.font.color.rgb = GRAY

prs.save(OUT)
print("저장:", OUT, "| 슬라이드", len(prs.slides.__iter__.__self__._sldIdLst))
